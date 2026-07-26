"""Conversion helpers for PDF <-> other formats.

These functions are intentionally simple — they prefer doing a faithful job
on common cases (text PDFs → DOCX, XLSX, PPTX) over trying to perfectly
reconstruct complex layouts.
"""
from __future__ import annotations

import os
import subprocess
import sys
from pathlib import Path

import fitz  # PyMuPDF

from .pdf_engine import EngineError, PAGE_SIZES


# ---------------------------------------------------------------------------
# PDF -> images
# ---------------------------------------------------------------------------

def pdf_to_images(src: str, dest_dir: str | os.PathLike,
                  fmt: str = "png", dpi: int = 200) -> list[str]:
    """Render every page to ``fmt`` ('png' | 'jpg' | 'tiff') at given dpi."""
    dest_dir = Path(dest_dir)
    dest_dir.mkdir(parents=True, exist_ok=True)
    fmt = fmt.lower()
    if fmt == "jpg":
        fmt = "jpeg"
    outputs: list[str] = []
    try:
        doc = fitz.open(src)
    except Exception as e:
        raise EngineError(f"Could not open PDF: {e}") from e
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    with doc:
        for i, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            dest = dest_dir / f"page_{i:03d}.{fmt}"
            pix.save(str(dest))
            outputs.append(str(dest))
    return outputs


# ---------------------------------------------------------------------------
# PDF -> text
# ---------------------------------------------------------------------------

def pdf_to_text(src: str, dest: str | os.PathLike) -> None:
    try:
        doc = fitz.open(src)
    except Exception as e:
        raise EngineError(f"Could not open PDF: {e}") from e
    chunks = []
    with doc:
        for i, page in enumerate(doc, start=1):
            chunks.append(f"\n\n----- Page {i} -----\n")
            chunks.append(page.get_text("text"))
    Path(dest).write_text("".join(chunks), encoding="utf-8")


# ---------------------------------------------------------------------------
# PDF -> Word (DOCX)
# ---------------------------------------------------------------------------

def pdf_to_word(src: str, dest: str | os.PathLike) -> None:
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError as e:
        raise EngineError("python-docx is required for PDF → Word.") from e
    try:
        doc = fitz.open(src)
    except Exception as e:
        raise EngineError(f"Could not open PDF: {e}") from e
    out = Document()
    with doc:
        for i, page in enumerate(doc, start=1):
            if i > 1:
                out.add_page_break()
            for line in page.get_text("text").splitlines():
                p = out.add_paragraph()
                run = p.add_run(line)
                run.font.size = Pt(11)
    out.save(str(dest))


# ---------------------------------------------------------------------------
# PDF -> Excel (XLSX)
# ---------------------------------------------------------------------------

def pdf_to_excel(src: str, dest: str | os.PathLike) -> None:
    try:
        from openpyxl import Workbook
    except ImportError as e:
        raise EngineError("openpyxl is required for PDF → Excel.") from e
    try:
        doc = fitz.open(src)
    except Exception as e:
        raise EngineError(f"Could not open PDF: {e}") from e
    out = Workbook()
    # Remove default sheet
    out.remove(out.active)
    with doc:
        # Try to use pdfplumber if available for true table extraction
        try:
            import pdfplumber
            with pdfplumber.open(src) as pdf:
                for i, page in enumerate(pdf.pages, start=1):
                    tables = page.extract_tables() or []
                    if not tables:
                        # Fallback: dump lines
                        sheet = out.create_sheet(f"Page {i}")
                        for r, line in enumerate((page.extract_text() or "").splitlines(), start=1):
                            sheet.cell(row=r, column=1, value=line)
                        continue
                    for t_idx, table in enumerate(tables, start=1):
                        sheet = out.create_sheet(f"Page {i} Table {t_idx}")
                        for r, row in enumerate(table, start=1):
                            for c, val in enumerate(row, start=1):
                                sheet.cell(row=r, column=c, value=(val or "").strip())
        except ImportError:
            for i, page in enumerate(doc, start=1):
                sheet = out.create_sheet(f"Page {i}")
                for r, line in enumerate(page.get_text("text").splitlines(), start=1):
                    sheet.cell(row=r, column=1, value=line)
    out.save(str(dest))


# ---------------------------------------------------------------------------
# PDF -> PowerPoint (PPTX)
# ---------------------------------------------------------------------------

def pdf_to_pptx(src: str, dest: str | os.PathLike) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:
        raise EngineError("python-pptx is required for PDF → PPT.") from e
    try:
        doc = fitz.open(src)
    except Exception as e:
        raise EngineError(f"Could not open PDF: {e}") from e

    pres = Presentation()
    blank_layout = pres.slide_layouts[6]
    with doc:
        # Use 16:9 widescreen
        pres.slide_width  = Inches(13.333)
        pres.slide_height = Inches(7.5)

        for i, page in enumerate(doc, start=1):
            slide = pres.slides.add_slide(blank_layout)
            mat = fitz.Matrix(1.5, 1.5)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            tmp_png = Path(dest).parent / f"_tmp_pptx_{i}.png"
            pix.save(str(tmp_png))
            slide.shapes.add_picture(
                str(tmp_png), 0, 0,
                width=pres.slide_width, height=pres.slide_height,
            )
            tmp_png.unlink(missing_ok=True)
    pres.save(str(dest))


# ---------------------------------------------------------------------------
# JPG / images -> PDF
# ---------------------------------------------------------------------------

def images_to_pdf(image_paths: list[str], dest: str | os.PathLike,
                  page_size: str = "A4") -> None:
    try:
        from PIL import Image as _Im
    except ImportError as e:
        raise EngineError("Pillow is required for images → PDF.") from e
    out = fitz.open()
    page_w, page_h = PAGE_SIZES.get(page_size, PAGE_SIZES["A4"])
    for p in image_paths:
        try:
            img = _Im.open(p)
        except Exception as e:
            raise EngineError(f"Could not open image {p}: {e}") from e
        # Convert to RGB if needed
        if img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGB")
        # Scale to fit while preserving aspect
        iw, ih = img.size
        scale = min(page_w / iw, page_h / ih) * 0.9
        new_w, new_h = iw * scale, ih * scale
        x = (page_w - new_w) / 2
        y = (page_h - new_h) / 2
        # Render onto a PDF page via Pixmap
        page = out.new_page(width=page_w, height=page_h)
        page.insert_image(
            fitz.Rect(x, y, x + new_w, y + new_h),
            filename=p,
        )
    out.save(str(dest), deflate=True, garbage=4)


# ---------------------------------------------------------------------------
# HTML -> PDF
# ---------------------------------------------------------------------------

def html_to_pdf(html: str, dest: str | os.PathLike) -> None:
    """Render a string of HTML to a PDF using WeasyPrint."""
    try:
        from weasyprint import HTML
    except (ImportError, OSError) as e:
        raise EngineError(
            "WeasyPrint requires native libraries (libgobject, libpango, etc.). "
            "On macOS install with: brew install cairo pango gdk-pixbuf libffi."
        ) from e
    try:
        HTML(string=html).write_pdf(str(dest))
    except Exception as e:
        raise EngineError(f"HTML rendering failed: {e}") from e


# ---------------------------------------------------------------------------
# Word (DOCX) -> PDF
# ---------------------------------------------------------------------------

def word_to_pdf(src: str, dest: str | os.PathLike) -> None:
    """Convert DOCX → PDF. On macOS we use Pages via AppleScript if available,
    otherwise a generic conversion that prints the text into a PDF.
    """
    src = str(src)
    dest = str(dest)
    # 1) Try Pages via AppleScript
    if sys.platform == "darwin":
        script = (
            f'tell application "Pages"\n'
            f'  set theDoc to open POSIX file "{src}"\n'
            f'  set thePath to POSIX file "{dest}"\n'
            f'  export theDoc to thePath as PDF\n'
            f'  close theDoc\n'
            f'end tell\n'
        )
        try:
            subprocess.run(["osascript", "-e", script], check=True, timeout=120)
            if os.path.exists(dest):
                return
        except Exception:
            pass  # fall through to plain text fallback

    # 2) Fallback: extract text from the docx and place in PDF
    try:
        from docx import Document
    except ImportError as e:
        raise EngineError("python-docx is required for Word → PDF.") from e
    doc = Document(src)
    text = "\n".join(p.text for p in doc.paragraphs)
    html = "<html><body><pre style='font-family: sans-serif'>" + \
           text.replace("&", "&amp;").replace("<", "&lt;") + \
           "</pre></body></html>"
    html_to_pdf(html, dest)
